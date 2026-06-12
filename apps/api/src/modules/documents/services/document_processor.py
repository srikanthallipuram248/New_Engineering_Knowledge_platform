from pathlib import Path
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import subprocess

from langchain_text_splitters import (
    RecursiveCharacterTextSplitter
)


class DocumentProcessor:

    @staticmethod
    def extract_text(file_path: str) -> str:
        path = Path(file_path)
        suffix = path.suffix.lower()

        if suffix == ".txt":
            return DocumentProcessor._extract_txt(path)

        elif suffix == ".pdf":
            return DocumentProcessor._extract_pdf(path)

        elif suffix == ".doc":
            return DocumentProcessor._extract_doc(path)

        elif suffix == ".docx":
            return DocumentProcessor._extract_docx(path)

        elif suffix == ".csv":
            return DocumentProcessor._extract_csv(path)

        elif suffix == ".xlsx":
            return DocumentProcessor._extract_xlsx(path)

        elif suffix == ".pptx":
            return DocumentProcessor._extract_pptx(path)

        else:
            raise ValueError(
                f"Unsupported file type: {suffix}"
            )

    @staticmethod
    def _extract_txt(path: Path) -> str:
        with open(
            path,
            "r",
            encoding="utf-8",
            errors="ignore"
        ) as file:
            return file.read()

    @staticmethod
    def _extract_pdf(path: Path) -> str:
        reader = PdfReader(str(path))

        text = ""

        for page in reader.pages:
            page_text = page.extract_text()

            if page_text:
                text += page_text + "\n"

        return text

    @staticmethod
    def _extract_doc(path: Path) -> str:

        try:
            result = subprocess.run(
                ["antiword", str(path)],
                capture_output=True,
                text=True
            )

            if result.returncode != 0:
                raise ValueError(result.stderr)

            return result.stdout

        except Exception as e:
            raise ValueError(
                f"Unable to read DOC file: {e}"
            )

    @staticmethod
    def _extract_docx(path: Path) -> str:
        doc = Document(str(path))

        return "\n".join(
            para.text.strip()
            for para in doc.paragraphs
            if para.text.strip()
        )

    @staticmethod
    def _extract_csv(path: Path) -> str:
        df = pd.read_csv(path)

        return df.to_string(index=False)

    @staticmethod
    def _extract_xlsx(path: Path) -> str:
        excel_file = pd.ExcelFile(path)

        text = ""

        for sheet in excel_file.sheet_names:

            df = pd.read_excel(
                path,
                sheet_name=sheet
            )

            text += f"\nSheet: {sheet}\n"
            text += df.to_string(
                index=False
            )
            text += "\n"

        return text

    @staticmethod
    def _extract_pptx(path: Path) -> str:
        presentation = Presentation(
            str(path)
        )

        text = []

        for slide in presentation.slides:
            for shape in slide.shapes:

                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text.append(
                            shape.text.strip()
                        )

        return "\n".join(text)

    @staticmethod
    def chunk_text(
        text: str,
        #chunk_size: int = 1000,
        #overlap: int = 100
    ) -> list[str]:

        if not text.strip():
            return []

        # chunks = []
        # start = 0

        # while start < len(text):

        #     end = start + chunk_size

        #     chunks.append(
        #         text[start:end]
        #     )

        #     start += (
        #         chunk_size - overlap
        #     )

        # return chunks

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=[
                "\n\n",
                "\n",
                ". ",
                " "
            ]
        )

        chunks = splitter.split_text(text)

        return chunks